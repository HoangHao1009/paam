import pandas as pd
from io import BytesIO
from typing import List, Union, Literal
from ..question import Single, Multiple, Number, Rank

from pptx import Presentation
from pptx.dml.color import MSO_THEME_COLOR, RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt

question_type = Union[Single, Multiple, Number, Rank]    

def _check_elements(list_check: List[str], questions: List[question_type]):
    question_codes = [question.code for question in questions]
    for question in list_check:
        if question not in question_codes:
            raise KeyError(f"Question {question} not exist in survey")

def _to_excel_buffer(df: pd.DataFrame, sheet_name: str, index=True) -> BytesIO:
    """Helper function to write a DataFrame to an Excel buffer."""
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=index, sheet_name=sheet_name)
    output.seek(0)
    return output

class PPTConfig:
    theme_color: List = [
    MSO_THEME_COLOR.ACCENT_1,
    MSO_THEME_COLOR.ACCENT_2,
    MSO_THEME_COLOR.ACCENT_3,
    MSO_THEME_COLOR.ACCENT_4,
    MSO_THEME_COLOR.ACCENT_5,
    MSO_THEME_COLOR.ACCENT_6,
]
    slide_layout: int = 5
    position: List[int] = [2, 2, 7, 4]
    font: str = 'Montserrat'
    has_legend: bool = True
    has_title: bool = True
    legend_font_size: int =  8
    category_axis_has_major_gridlines: bool = False
    category_axis_has_minor_gridlines: bool = False
    category_axis_has_title: bool = False
    category_axis_visible: bool = True
    category_axis_tick_labels_font_size: int = 12
    value_axis_has_major_gridlines: bool = False
    value_axis_has_minor_gridlines: bool = False
    value_axis_visible: bool = False
    data_labels_font_size: int = 8
    data_labels_font: int = 'Montserrat'
    data_labels_number_format: int = 'General'
    data_labels_number_format_is_linked: bool = True
    _data_labels_position: Literal['above', 'below', 'best_fit', 'center', 'inside_base', 'inside_end', 'left', 'mixed', 'outside_end', 'right'] = 'outside_end'
    data_labels_show_category_name: bool = False
    data_labels_show_legend_key: bool = False
    data_labels_show_percentage: bool = False
    data_labels_show_series_name: bool = False
    data_labels_show_value: bool = True

def _add_chart_to_prs(
        prs: Presentation, df: pd.DataFrame, 
        title, config: PPTConfig, 
        from_template=False,
        question_text_1: str="",
        question_text_2: str=""
    ):
    
    def emu_to_inches(emu):
        return Inches(emu/ 914400)
    
    def add_textbox(text, x, y, width, height, font_name, font_size=13, italic=True):
        text_box = slide.shapes.add_textbox(x, y, width, height)
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.italic = italic
        text_frame.word_wrap = False
        
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    
    slide_layout = prs.slides[0].slide_layout if from_template else prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    chart_data = CategoryChartData()
    chart_data.categories = [str(cat) for cat in df.index.get_level_values(0)]
    for col in df.columns:
        chart_data.add_series(str(col[0]), df[col])
        
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    chart_width = slide_width / 2
    chart_height = slide_height / 2
    x = emu_to_inches((slide_width - chart_width) / 2)
    y = emu_to_inches((slide_height - chart_height) / 2)
    cx = emu_to_inches(chart_width)
    cy = emu_to_inches(chart_height)
            
    chart = slide.shapes.add_chart(
        chart_type=chart_type,
        x=x, y=y, cx=cx, cy=cy,
        chart_data=chart_data
    ).chart
    
    chart.font.name = config.font
    chart.has_legend = config.has_legend
    chart.has_title = config.has_title
    chart.chart_title.text_frame.text = title
    
    title_font = chart.chart_title.text_frame.paragraphs[0].font
    title_font.name = config.font
    title_font.size = Pt(14)

    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(config.legend_font_size)
    try:
        chart.category_axis.has_major_gridlines = config.category_axis_has_major_gridlines
        chart.category_axis.has_minor_gridlines = config.category_axis_has_minor_gridlines
        chart.category_axis.has_title = config.category_axis_has_title
        chart.category_axis.visible = config.category_axis_visible
        chart.category_axis.tick_labels.font.size = Pt(config.category_axis_tick_labels_font_size)
        chart.value_axis.has_major_gridlines = config.value_axis_has_major_gridlines
        chart.value_axis.has_minor_gridlines = config.value_axis_has_minor_gridlines
        chart.value_axis.visible = config.value_axis_visible
    except:
        pass

    fill = chart.plots[0].series[0].format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(79, 129, 189)
    
    text_width = Inches(2)
    text_height = Inches(0.5)
    text_x = Inches(0.2)
    text_y1 = Inches(slide_height / 914400 - 0.9)
    text_y2 = Inches(slide_height / 914400 - 0.7)
    
    add_textbox(question_text_1, text_x, text_y1, text_width, text_height, config.font, 10, True)
    add_textbox(question_text_2, text_x, text_y2, text_width, text_height, config.font, 10, True)
    
