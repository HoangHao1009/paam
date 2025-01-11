import pandas as pd
from typing import Union, List, Dict
from io import BytesIO
import pyreadstat
import zipfile
import tempfile

from ..question import Single, Multiple, Number, Rank, Answer
from ..ctab import CrossTab
from .utils import _to_excel_buffer, _check_elements, _add_chart_to_prs, PPTConfig
from ..utils import spss_function
from ..config import SPSSConfig, DataFrameConfig

from pptx import Presentation
import ast

question_type = Union[Single, Multiple, Number, Rank]

class Survey:
    def __init__(self, data: dict, 
                 control_vars: List[Union[Single, Multiple]]=[],
                 target_vars: List[question_type]=[], 
                 deep_vars: List[Union[Single, Multiple]]=[]
            ):
        self.data = data
        self._df_config = DataFrameConfig()
        self._spss_config = SPSSConfig()
        self.control_vars = control_vars
        self.target_vars = target_vars
        self.deep_vars = deep_vars
        
    def initialize(self):
        self.questions: List[question_type] = self._get_questions()
        self._check_valid()
        
    @property
    def df_config(self):
        return self._df_config
    
    @property
    def spss_config(self):
        return self._spss_config
    
    @property
    def stats(self):
        return {
            'num_questions': len(self.data['questions']),
            'num_respondents': len(self.data['respondents'])
        }
        
    @property
    def question_code_mapping(self):
        return {question.code: question for question in self.questions}
        
    def _check_valid(self):
        list_question_code = [question.code for question in self.questions]
        if len(list_question_code) != len(set(list_question_code)):
            seen = set()
            duplicates = set()
            for i in list_question_code:
                if i in seen:
                    duplicates.add(i)
                else:
                    seen.add(i)
            raise KeyError(f'Question code: {list(duplicates)} duplicated')
                
    def set_df_config(self, settings: dict):
        """
        Set config for .df
        Example:
        {'scale_encode': True}
        """
        for k, v in settings.items():
            if hasattr(self._df_config, k):
                setattr(self._df_config, k, v)
        for question in self.questions:
            for k, v in settings.items():
                if hasattr(question, k):
                    setattr(question, k, v)
                    
    def set_spss_config(self, settings: dict):
        """
        Set config for .df
        Example:
        {'std': True, 'pct': True, 'compare_tests': ['MEAN', 'PROP'], 'alpha': 0.1}
        """
        for k, v in settings.items():
            if hasattr(self._spss_config, k):
                setattr(self._spss_config, k, v)
                
    def _get_questions(self):
        answer_lookup = {answer['id']: answer for answer in self.data['answers']}
        all_questions = []
        for question_data in self.data['questions']:
            answers = []
            for answer_id in question_data['answers']:
                answer_data = answer_lookup[answer_id]
                answer = Answer(**answer_data, question_code=question_data['code'])
                answer.is_rank = True if question_data['type'] == 'rank' else False
                answers.append(answer)
                
            question_data['answers'] = answers
                
            if question_data['type'] in ['sa', 'sa_matrix', 'text']:
                question = Single(**question_data)
            elif question_data['type'] in ['ma', 'ma_matrix']:
                question = Multiple(**question_data)
            elif question_data['type'] in ['number']:
                question = Number(**question_data)
            elif question_data['type'] in ['rank']:
                question = Rank(**question_data)
            else:
                raise ValueError(f"Question with id: {question_data['id']} code: {question_data['code']} with type: {question_data['type']} can not be processed")
            
            question.sort_answer(list(range(1, len(question.answers) + 1)), by='scale')
            
            all_questions.append(question)
            all_questions.sort(key=lambda item: item.order) 
        
        return all_questions
  
    def _get_constructed_questions(self, compute_constructions: List[Dict]):
        for construction in compute_constructions:
            root_question = self.get_question(construction['rootQuestionCode'])
            new_question = root_question.compute(
                construct_dict=ast.literal_eval(construction['construction']),
                new_code=construction['newQuestionCode'],
                new_text=construction['newQuestionText'],
                method=construction['method'],
                by=construction['by'],
            )
            self.questions.append(new_question)
            self.questions.sort(key=lambda item: item.order)
             
    def get_question(self, key: Union[List[str], str]):
        """
        Getting questions by question code
        Arguments:
        - key: list or str of question code
        Return:
        List or single question object
        """
        if isinstance(key, str):
            if key in self.question_code_mapping:
                return self.question_code_mapping[key]
            else:
                raise KeyError(f"Question code {key} cannot be found")
        else:
            questions = []
            for question_code in key:
                if question_code in self.question_code_mapping:
                    questions.append(self.question_code_mapping[question_code])
                else:
                    raise KeyError(f"Question code {question_code} cannot be found")
            return questions
        
    def __getitem__(self, key: Union[List[str], str]):
        return self.get_question(key)
    
    def add_question(self, question_obj: question_type):
        if question_obj.code in self.question_code_mapping:
            raise KeyError(f"Question code: {question_obj} has been existed in survey")
        question_order = [question.order for question in self.questions]
        if question_obj.order in question_order:
            raise ValueError(f"Question code {question_obj.code} order: {question_obj.order} has been existed in survey")
        self.questions.append(question_obj)
        
    @property
    def df(self):
        for question in self.questions:
            question._ctab_mode = False
        dfs = [question.df for question in self.questions]
        return pd.concat(dfs, axis=1).reset_index()
        
    def crosstab(self, base, target, deep_by: list=[], alpha: float=0, pct: bool=False) -> CrossTab:
        """
        Create a crosstab for 2 questions
        Parameters:
        base: variable in column
        target: varibale in row
        deep_by: list of filter variable
        """
        base_question = self.get_question(base)
        target_question = self.get_question(target)
        deep_by = [self.get_question(deep) for deep in deep_by]
        ctab = CrossTab(base=base_question, target=target_question, deep_by=deep_by)
        ctab.config.alpha = alpha
        ctab.config.pct = pct
        return ctab
            
    @property
    def ctab_df(self):
        if self.df_config.scale_encode:
            self.set_df_config({'scale_encode': True})
        
        if not hasattr(self, 'target_vars'):
            self.target_vars = [question.code for question in self.questions]

        if not hasattr(self, 'deep_vars'):
            self.deep_vars = []
        
        if hasattr(self, 'control_vars'):
            deep_by = [self.get_question(var) for var in self.deep_vars]
            
            dfs = []
            for row_var in self.target_vars:
                cols = []
                for col_var in self.control_vars:
                    base = self.get_question(col_var)
                    target = self.get_question(row_var)
                    try:
                        col_ctab_df = CrossTab(
                            base=base,
                            target=target,
                            deep_by=deep_by
                        ).df
                    except Exception as e:
                        print(f"Can not crosstab with base: {col_var} - target: {row_var} -> {e}")
                        
                    col_ctab_df.index = pd.MultiIndex.from_product([[col_ctab_df.index.name], col_ctab_df.index.to_list()])
                    
                    cols.append(col_ctab_df)
                dfs.append(pd.concat(cols, axis=1))
            return pd.concat(dfs, axis=0)  
        else:
            raise ValueError(f"Required control variables to calculate crosstab")

    def to_excel(self) -> tuple[BytesIO, BytesIO]:
        """Export both raw data and cross-tab data to Excel buffers."""
        if self.df.empty:
            raise ValueError("Raw data DataFrame is empty")
        if self.ctab_df.empty:
            raise ValueError("CrossTab DataFrame is empty")

        rawdata_buffer = _to_excel_buffer(self.df, 'RawData')
        ctab_buffer = _to_excel_buffer(self.ctab_df, 'CrossTab')
        
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            # Ghi nội dung từ buffer vào file ZIP
            zip_file.writestr("RawData.xlsx", rawdata_buffer.getvalue())
            zip_file.writestr("CrossTab.xlsx", ctab_buffer.getvalue())
        zip_buffer.seek(0)

        return zip_buffer
    
    @property
    def spss_syntax(self):
        syntax = []
        for question in self.questions:
            syntax.extend(question.spss_syntax)
        
        if self.control_vars:
            calculate_dict = {}
            for code in [question.code for question in self.questions]:
                question = self[code]
                                        
                if isinstance(question, Number):
                    calculate_dict[code] = ['Mean', 'Std'] if self.spss_config.std else ['Mean']
                else:
                    if isinstance(question, Multiple):
                        calculate_dict[f'${code}'] = ['ColPct'] if self.spss_config.pct else ['Count']
                    else:
                        calculate_dict[code] = ['ColPct'] if self.spss_config.pct else ['Count']
            by_col = []
            for var in self.control_vars:
                var = f'${var}' if isinstance(self[var], Multiple) else var
                by_col.append(var)
            ctab_syntax = spss_function.ctab(by_col, calculate_dict, self.spss_config.compare_tests, self.spss_config.alpha)
            syntax.append(ctab_syntax)

        return '\n'.join(syntax)
    
    def to_spss(self) -> tuple[BytesIO, BytesIO]:
        """Export both sav and sps file to buffers."""
        df_old_config = self.df_config.__dict__
        self.set_df_config({'scale_encode': True})
        
        df = self.df
        
        with tempfile.NamedTemporaryFile(suffix='.sav', delete=False) as temp_sav_file:
            pyreadstat.write_sav(df, temp_sav_file.name)
            temp_sav_file.seek(0)
            sav_output = BytesIO(temp_sav_file.read())
            sav_output.seek(0)
        
        self.set_df_config(df_old_config)
        
        sps_output = BytesIO()
        sps_output.write(self.spss_syntax.encode('utf-8'))
        sps_output.seek(0)
        
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            zip_file.writestr("data.sav", sav_output.getvalue())
            zip_file.writestr("syntax.sps", sps_output.getvalue())
        
        zip_buffer.seek(0)
        
        return zip_buffer
    
    def to_datasets(self):
        question_data = [
            {
                "question_id": question._id,
                "question_code": question.code, 
                "question_text": question.text, 
                "question_type": question.type,
                "question_order": question.order,
            }
            for question in self.questions
        ]
        
        dim_question = pd.DataFrame(question_data)
        
        answer_data = [
            {
                "answer_id": answer._id,
                "answer_code": answer.code, 
                "answer_text": answer.text, 
                "answer_scale": answer.scale,
                "question_code": answer.question_code
            }
            for question in self.questions for answer in question.answers
        ]
        
        dim_answer = pd.DataFrame(answer_data)
                
        fact_data = [
            {
                "respondent_id": respondent,
                "question_id": question._id,
                "question_code": question.code,
                "answer_id": answer._id,
                "answer_code": answer.code,
                "answer_text": answer.text,
                "answer_scale": answer.scale
            }
            for question in self.questions for answer in question.answers for respondent in answer.respondents
        ]
        
        fact = pd.DataFrame(fact_data)
        
        self.set_df_config({'scale_encode': False})
        dim_respondent_intext = self.df
        self.set_df_config({'scale_encode': True})
        dim_respondent_inscale = self.df
        
        dim_question_buffer = _to_excel_buffer(dim_question, 'dimQuestion', index=False)
        dim_answer_buffer = _to_excel_buffer(dim_answer, 'dimAnswer', index=False)
        fact_buffer = _to_excel_buffer(fact, 'fact', index=False)
        dim_respondent_intext_buffer = _to_excel_buffer(dim_respondent_intext, 'dimRespondentInText', index=True)
        dim_respondent_inscale_buffer = _to_excel_buffer(dim_respondent_inscale, 'dimRespondentInScale', index=True)
        
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            # Ghi nội dung từ buffer vào file ZIP
            zip_file.writestr("dimQuestion.xlsx", dim_question_buffer.getvalue())
            zip_file.writestr("dimAnswer.xlsx", dim_answer_buffer.getvalue())
            zip_file.writestr("fact.xlsx", fact_buffer.getvalue())
            zip_file.writestr("dimRespondentInText.xlsx", dim_respondent_intext_buffer.getvalue())
            zip_file.writestr("dimRespondentInScale.xlsx", dim_respondent_inscale_buffer.getvalue())

        zip_buffer.seek(0)

        return zip_buffer
    
    def to_ppt(self, config: PPTConfig=PPTConfig(), template_path: str=""):
        prs = Presentation(template_path) if template_path else Presentation()
        from_template = True if template_path else False
        for base in self.control_vars:
            for target in self.target_vars:
                ctab = self.crosstab(base, target)
                # ctab.config.total = False
                df = ctab.df
                _add_chart_to_prs(prs, df, f'{base}_{target}', config=config, from_template=from_template)
                
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            zip_file.writestr("PPTX_Chart.pptx", pptx_buffer.read())

        zip_buffer.seek(0)
        
        return zip_buffer


